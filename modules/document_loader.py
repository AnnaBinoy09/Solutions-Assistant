"""
modules/document_loader.py — Extended Document Loader
───────────────────────────────────────────────────────
Supports:
  .pdf    — PyPDF page-by-page extraction
  .docx   — python-docx paragraph grouping
  .txt    — plain text, encoding-safe
  .md     — markdown stripped to plain text
  .csv    — pandas tabular → text rows
  .xlsx   — openpyxl sheet-by-sheet extraction
  .pptx   — python-pptx slide-by-slide text
  .html   — BeautifulSoup tag stripping
  .png /
  .jpg /
  .jpeg /
  .tiff /
  .bmp    — pytesseract OCR
"""

import os
import logging
from typing import List
from dataclasses import dataclass, field

logger = logging.getLogger(__name__)


@dataclass
class Document:
    page_content: str
    metadata: dict = field(default_factory=dict)

    def __repr__(self):
        snippet = self.page_content[:80].replace("\n", " ")
        return (
            f"Document(source={self.metadata.get('source', '?')!r}, "
            f"page={self.metadata.get('page', '?')}, content={snippet!r}...)"
        )


class DocumentLoader:
    """
    Unified document loader.  One public method — load() — dispatches
    to the correct private loader based on file extension.
    """

    SUPPORTED_EXTENSIONS = {
        ".pdf", ".docx",
        ".txt", ".md",
        ".csv", ".xlsx",
        ".pptx",
        ".html", ".htm",
        ".png", ".jpg", ".jpeg", ".tiff", ".bmp",
    }

    # ──────────────────────────────────────────
    # Public API
    # ──────────────────────────────────────────

    def load(self, file_path: str) -> List[Document]:
        """
        Load any supported file and return a list of Document objects.

        Args:
            file_path: Absolute or relative path to the document.

        Returns:
            List[Document] — each with page_content and metadata.

        Raises:
            FileNotFoundError: If the file does not exist.
            ValueError: If the file type is unsupported.
        """
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"File not found: {file_path}")

        ext = os.path.splitext(file_path)[-1].lower()
        if ext not in self.SUPPORTED_EXTENSIONS:
            raise ValueError(
                f"Unsupported file type '{ext}'. "
                f"Supported: {sorted(self.SUPPORTED_EXTENSIONS)}"
            )

        logger.info(f"Loading document: {file_path} (type={ext})")

        dispatch = {
            ".pdf":   self._load_pdf,
            ".docx":  self._load_docx,
            ".txt":   self._load_txt,
            ".md":    self._load_md,
            ".csv":   self._load_csv,
            ".xlsx":  self._load_xlsx,
            ".pptx":  self._load_pptx,
            ".html":  self._load_html,
            ".htm":   self._load_html,
            ".png":   self._load_image,
            ".jpg":   self._load_image,
            ".jpeg":  self._load_image,
            ".tiff":  self._load_image,
            ".bmp":   self._load_image,
        }
        return dispatch[ext](file_path)

    # ──────────────────────────────────────────
    # PDF
    # ──────────────────────────────────────────

    def _load_pdf(self, file_path: str) -> List[Document]:
        try:
            import pypdf
        except ImportError:
            raise ImportError("pip install pypdf")

        documents = []
        source_name = os.path.basename(file_path)

        with open(file_path, "rb") as f:
            reader = pypdf.PdfReader(f)
            total_pages = len(reader.pages)
            logger.info(f"  PDF has {total_pages} pages.")

            for page_num, page in enumerate(reader.pages, start=1):
                text = (page.extract_text() or "").strip()
                if not text:
                    continue
                documents.append(Document(
                    page_content=text,
                    metadata={
                        "source": source_name,
                        "file_path": file_path,
                        "page": page_num,
                        "total_pages": total_pages,
                        "file_type": "pdf",
                    }
                ))

        logger.info(f"  Loaded {len(documents)} non-empty pages from PDF.")
        return documents

    # ──────────────────────────────────────────
    # DOCX
    # ──────────────────────────────────────────

    def _load_docx(self, file_path: str) -> List[Document]:
        try:
            from docx import Document as DocxDocument
        except ImportError:
            raise ImportError("pip install python-docx")

        doc = DocxDocument(file_path)
        source_name = os.path.basename(file_path)
        paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]

        if not paragraphs:
            logger.warning(f"DOCX appears empty: {file_path}")
            return []

        section_size = 30
        sections = [
            paragraphs[i: i + section_size]
            for i in range(0, len(paragraphs), section_size)
        ]
        total_sections = len(sections)
        documents = []
        for sec_num, sec_paras in enumerate(sections, start=1):
            documents.append(Document(
                page_content="\n\n".join(sec_paras),
                metadata={
                    "source": source_name,
                    "file_path": file_path,
                    "page": sec_num,
                    "total_pages": total_sections,
                    "file_type": "docx",
                }
            ))

        logger.info(f"  Loaded {len(documents)} sections from DOCX.")
        return documents

    # ──────────────────────────────────────────
    # TXT
    # ──────────────────────────────────────────

    def _load_txt(self, file_path: str) -> List[Document]:
        source_name = os.path.basename(file_path)
        for enc in ("utf-8", "latin-1", "cp1252"):
            try:
                with open(file_path, "r", encoding=enc) as f:
                    text = f.read().strip()
                break
            except UnicodeDecodeError:
                continue
        else:
            raise ValueError(f"Cannot decode text file: {file_path}")

        if not text:
            logger.warning(f"TXT file is empty: {file_path}")
            return []

        # Split into paragraphs; group into page-like sections
        paragraphs = [p.strip() for p in text.split("\n\n") if p.strip()]
        section_size = 20
        sections = [
            paragraphs[i: i + section_size]
            for i in range(0, len(paragraphs), section_size)
        ]
        documents = []
        for sec_num, sec_paras in enumerate(sections, start=1):
            documents.append(Document(
                page_content="\n\n".join(sec_paras),
                metadata={
                    "source": source_name,
                    "file_path": file_path,
                    "page": sec_num,
                    "total_pages": len(sections),
                    "file_type": "txt",
                }
            ))

        logger.info(f"  Loaded {len(documents)} sections from TXT.")
        return documents

    # ──────────────────────────────────────────
    # Markdown
    # ──────────────────────────────────────────

    def _load_md(self, file_path: str) -> List[Document]:
        """Strip markdown syntax and treat as plain text sections."""
        source_name = os.path.basename(file_path)
        for enc in ("utf-8", "latin-1"):
            try:
                with open(file_path, "r", encoding=enc) as f:
                    raw = f.read()
                break
            except UnicodeDecodeError:
                continue
        else:
            raise ValueError(f"Cannot decode markdown file: {file_path}")

        # Try markdownify for clean plain-text, fall back to raw
        try:
            from markdownify import markdownify as md_to_text
            text = md_to_text(raw).strip()
        except ImportError:
            import re
            text = re.sub(r"[#*`>\-_~\[\]]+", " ", raw).strip()

        if not text:
            return []

        paragraphs = [p.strip() for p in text.split("\n\n") if p.strip()]
        section_size = 20
        sections = [
            paragraphs[i: i + section_size]
            for i in range(0, len(paragraphs), section_size)
        ]
        documents = []
        for sec_num, sec_paras in enumerate(sections, start=1):
            documents.append(Document(
                page_content="\n\n".join(sec_paras),
                metadata={
                    "source": source_name,
                    "file_path": file_path,
                    "page": sec_num,
                    "total_pages": len(sections),
                    "file_type": "md",
                }
            ))

        logger.info(f"  Loaded {len(documents)} sections from Markdown.")
        return documents

    # ──────────────────────────────────────────
    # CSV
    # ──────────────────────────────────────────

    def _load_csv(self, file_path: str) -> List[Document]:
        try:
            import pandas as pd
        except ImportError:
            raise ImportError("pip install pandas")

        source_name = os.path.basename(file_path)
        df = pd.read_csv(file_path, encoding_errors="replace")

        if df.empty:
            logger.warning(f"CSV is empty: {file_path}")
            return []

        # Convert schema summary as first document
        schema_text = (
            f"File: {source_name}\n"
            f"Columns ({len(df.columns)}): {', '.join(df.columns.tolist())}\n"
            f"Rows: {len(df)}\n\n"
            f"Sample (first 5 rows):\n{df.head(5).to_string(index=False)}"
        )

        documents = [Document(
            page_content=schema_text,
            metadata={
                "source": source_name,
                "file_path": file_path,
                "page": 1,
                "total_pages": 1,
                "file_type": "csv",
                "rows": len(df),
                "columns": len(df.columns),
            }
        )]

        # Batch remaining rows into pages of 50
        batch_size = 50
        for batch_num, start in enumerate(range(0, len(df), batch_size), start=2):
            batch = df.iloc[start: start + batch_size]
            documents.append(Document(
                page_content=batch.to_string(index=False),
                metadata={
                    "source": source_name,
                    "file_path": file_path,
                    "page": batch_num,
                    "total_pages": (len(df) // batch_size) + 2,
                    "file_type": "csv",
                }
            ))

        logger.info(f"  Loaded {len(documents)} sections from CSV.")
        return documents

    # ──────────────────────────────────────────
    # XLSX
    # ──────────────────────────────────────────

    def _load_xlsx(self, file_path: str) -> List[Document]:
        try:
            import pandas as pd
        except ImportError:
            raise ImportError("pip install pandas openpyxl")

        source_name = os.path.basename(file_path)
        xl = pd.ExcelFile(file_path, engine="openpyxl")
        documents = []
        page_num = 1

        for sheet_name in xl.sheet_names:
            df = xl.parse(sheet_name)
            if df.empty:
                continue

            schema_text = (
                f"Sheet: {sheet_name} | File: {source_name}\n"
                f"Columns ({len(df.columns)}): {', '.join(str(c) for c in df.columns)}\n"
                f"Rows: {len(df)}\n\n"
                f"Sample (first 5 rows):\n{df.head(5).to_string(index=False)}"
            )

            documents.append(Document(
                page_content=schema_text,
                metadata={
                    "source": source_name,
                    "file_path": file_path,
                    "page": page_num,
                    "total_pages": len(xl.sheet_names),
                    "file_type": "xlsx",
                    "sheet": sheet_name,
                }
            ))
            page_num += 1

            # Remaining rows in batches of 50
            batch_size = 50
            for start in range(0, len(df), batch_size):
                batch = df.iloc[start: start + batch_size]
                documents.append(Document(
                    page_content=f"Sheet: {sheet_name}\n{batch.to_string(index=False)}",
                    metadata={
                        "source": source_name,
                        "file_path": file_path,
                        "page": page_num,
                        "total_pages": len(xl.sheet_names),
                        "file_type": "xlsx",
                        "sheet": sheet_name,
                    }
                ))
                page_num += 1

        logger.info(f"  Loaded {len(documents)} sections from XLSX.")
        return documents

    # ──────────────────────────────────────────
    # PPTX
    # ──────────────────────────────────────────

    def _load_pptx(self, file_path: str) -> List[Document]:
        try:
            from pptx import Presentation
        except ImportError:
            raise ImportError("pip install python-pptx")

        source_name = os.path.basename(file_path)
        prs = Presentation(file_path)
        documents = []
        total_slides = len(prs.slides)

        for slide_num, slide in enumerate(prs.slides, start=1):
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text.strip())

            if not texts:
                continue

            slide_text = f"[Slide {slide_num}/{total_slides}]\n" + "\n\n".join(texts)
            documents.append(Document(
                page_content=slide_text,
                metadata={
                    "source": source_name,
                    "file_path": file_path,
                    "page": slide_num,
                    "total_pages": total_slides,
                    "file_type": "pptx",
                }
            ))

        logger.info(f"  Loaded {len(documents)} slides from PPTX.")
        return documents

    # ──────────────────────────────────────────
    # HTML
    # ──────────────────────────────────────────

    def _load_html(self, file_path: str) -> List[Document]:
        try:
            from bs4 import BeautifulSoup
        except ImportError:
            raise ImportError("pip install beautifulsoup4")

        source_name = os.path.basename(file_path)
        for enc in ("utf-8", "latin-1"):
            try:
                with open(file_path, "r", encoding=enc) as f:
                    raw = f.read()
                break
            except UnicodeDecodeError:
                continue
        else:
            raise ValueError(f"Cannot decode HTML file: {file_path}")

        soup = BeautifulSoup(raw, "html.parser")
        # Remove script and style tags
        for tag in soup(["script", "style", "nav", "footer", "head"]):
            tag.decompose()

        text = soup.get_text(separator="\n")
        # Collapse whitespace
        import re
        text = re.sub(r"\n{3,}", "\n\n", text).strip()

        if not text:
            return []

        paragraphs = [p.strip() for p in text.split("\n\n") if p.strip()]
        section_size = 20
        sections = [
            paragraphs[i: i + section_size]
            for i in range(0, len(paragraphs), section_size)
        ]
        documents = []
        for sec_num, sec_paras in enumerate(sections, start=1):
            documents.append(Document(
                page_content="\n\n".join(sec_paras),
                metadata={
                    "source": source_name,
                    "file_path": file_path,
                    "page": sec_num,
                    "total_pages": len(sections),
                    "file_type": "html",
                }
            ))

        logger.info(f"  Loaded {len(documents)} sections from HTML.")
        return documents

    # ──────────────────────────────────────────
    # Images (OCR)
    # ──────────────────────────────────────────

    def _load_image(self, file_path: str) -> List[Document]:
        try:
            import pytesseract
            from PIL import Image
        except ImportError:
            raise ImportError("pip install pytesseract pillow")

        # ── Windows: point pytesseract to the binary ──
        import sys
        if sys.platform == "win32":
            try:
                import pytesseract
                pytesseract.pytesseract.tesseract_cmd = r"D:\Tesseract-OCR\tesseract.exe"
            except ImportError:
                pass
            pytesseract.pytesseract.tesseract_cmd = (
                r"D:\Tesseract-OCR\tesseract.exe"
            )

        source_name = os.path.basename(file_path)
        img = Image.open(file_path)

        try:
            text = pytesseract.image_to_string(img).strip()
        except Exception as e:
            raise RuntimeError(
                f"Tesseract OCR failed on '{source_name}'. "
                "Ensure tesseract is installed on the system: "
                "https://github.com/tesseract-ocr/tesseract\n"
                f"Error: {e}"
            )

        if not text:
            logger.warning(f"OCR produced no text from image: {source_name}")
            return []

        documents = [Document(
            page_content=text,
            metadata={
                "source": source_name,
                "file_path": file_path,
                "page": 1,
                "total_pages": 1,
                "file_type": os.path.splitext(file_path)[-1].lstrip("."),
                "ocr": True,
            }
        )]

        logger.info(f"  OCR extracted {len(text)} chars from image.")
        return documents

    # ──────────────────────────────────────────
    # Utility
    # ──────────────────────────────────────────

    @staticmethod
    def is_supported(file_path: str) -> bool:
        ext = os.path.splitext(file_path)[-1].lower()
        return ext in DocumentLoader.SUPPORTED_EXTENSIONS